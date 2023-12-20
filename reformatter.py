# %% IMPORT THE NECESSARY MODULES ***********************************************
from pptx.oxml.xmlchemy import OxmlElement
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os
import argparse
import re

# %% DEFINE THE FOLDER TO FORMAT USING COMMAND LINE ARGUMENTS *******************
parser = argparse.ArgumentParser()
parser.add_argument(
    '-d', '--directory', default='./', help='Directory to format')
parser.add_argument(
    '-f', '--file', default='./', help='File to format')
args = parser.parse_args()
folderToFormat = args.directory
fileToFormat = args.file

# %% FORMATTING FUNCTION ********************************************************


def apply_format(run, format_spec):
    try:
        run.font.name = format_spec['font_name']
        # NOTE: This code has the capability to change size, bold, and italic
        # but it is not currently used. Uncomment the lines below to use it.
        # You will also neeed to update the keys in the words.json file.
        run.font.size = run.font.size
        # run.font.bold = format_spec['font_bold']
        # run.font.italic = format_spec['font_italic']

        run.font.color.rgb = RGBColor(*format_spec['font_color'])
    except TypeError:
        return


def format_run(paragraph, word, format_spec):
    new_runs = []  # List to store new runs

    for run in paragraph.runs:
        split_text = run.text.split(word)

        # Check if the word is part of another word
        try:
            partOfWord = re.findall(r'\b' + word + r'\b', run.text.lower())
            partOfWord = False if partOfWord else True

        except:
            partOfWord = False

        if word in run.text.lower() and partOfWord == False:

            # Iterate through each part in split_text
            for i, part in enumerate(split_text):
                # Add the non-word part
                if part:
                    new_run = paragraph.add_run()
                    new_run.text = part
                    new_run = _copy_run_formatting(new_run, run)
                    new_runs.append(new_run)

                # Add the formatted word, except for the last part if it's empty
                if i < len(split_text) - 1:
                    formatted_run = paragraph.add_run()
                    formatted_run.text = word
                    apply_format(formatted_run, format_spec)
                    new_runs.append(formatted_run)
        elif run.text != '':
            # If the word is not in the run, just copy the run as is
            new_run = paragraph.add_run()
            new_run.text = run.text
            new_run = _copy_run_formatting(new_run, run)
            new_runs.append(new_run)

    # Clear original runs and update paragraph with new runs
    paragraph.clear()
    for new_run in new_runs:
        # Add the new run to the paragraph's xml
        run = paragraph.add_run()
        run.text = new_run.text
        run = _copy_run_formatting(run, new_run)

    # Clear the temporary new_runs list

    return paragraph


def _copy_run_formatting(new_run, original_run):
    # Copy formatting from original run to new run
    try:
        new_run.font.name = original_run.font.name
        new_run.font.size = original_run.font.size
        new_run.font.bold = original_run.font.bold
        new_run.font.italic = original_run.font.italic
        new_run.font.color.rgb = original_run.font.color.rgb
        return new_run
    except:
        return new_run


# %% MAIN PROGRAM ***************************************************************
with open('words.json', 'r') as file:
    formatting_rules = json.load(file)

if not os.path.exists(folderToFormat):
    print('Folder does not exist')
    exit()

# file parsing logic
if folderToFormat == './' and fileToFormat != './':
    files = [os.path.basename(fileToFormat)]
    folderToFormat = os.path.dirname(fileToFormat)
elif folderToFormat != './' and fileToFormat == './':
    files = os.listdir(folderToFormat)
    files = [f for f in files if f.endswith('.pptx')]
elif folderToFormat == './' and fileToFormat == './':
    files = os.listdir(folderToFormat)
    files = [f for f in files if f.endswith('.pptx')]


for fileToFormat in files:
    pptx = Presentation(folderToFormat+'/'+fileToFormat)
    for slide in pptx.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for rule in formatting_rules['words']:
                        # Iterate over each word in the "words" list
                        for word in rule['words']:
                            if word in paragraph.text.lower():
                                paragraph = format_run(
                                    paragraph, word, rule['format'])

    # Save the formatted presentation
    pptx.save(folderToFormat+'/'+fileToFormat.replace('.pptx', '.pptx'))
