import json
from pptx import Presentation
from pptx.dml.color import RGBColor


def apply_format(run, format_spec):
    run.font.name = format_spec.get('font_name', run.font.name)
    try:
        run.font.color.rgb = RGBColor(*format_spec['font_color'])
    except:
        return
    # Add more formatting options here as needed


def format_run(paragraph, word, format_spec):
    for run in paragraph.runs:
        if word in run.text:
            new_text_parts = []
            split_text = run.text.split(word)

            # Build new text with formatting
            for i, part in enumerate(split_text):
                new_text_parts.append(part)
                if i < len(split_text) - 1:
                    new_text_parts.append(word)

            # Replace run text and apply formatting
            run.text = new_text_parts[0]
            for text_part in new_text_parts[1:]:
                if text_part == word:
                    formatted_run = paragraph.add_run()
                    formatted_run.text = word
                    apply_format(formatted_run, format_spec)
                else:
                    new_run = paragraph.add_run()
                    new_run.text = text_part
                    _copy_run_formatting(new_run, run)


def _copy_run_formatting(new_run, original_run):
    new_run.font.name = original_run.font.name
    new_run.font.size = original_run.font.size
    new_run.font.bold = original_run.font.bold
    new_run.font.italic = original_run.font.italic
    try:
        new_run.font.color.rgb = original_run.font.color.rgb
    except:
        return


# Load the presentation and formatting rules
pptx = Presentation('06 theWhileLoop.pptx')
with open('words.json', 'r') as file:
    formatting_rules = json.load(file)

# Iterate through each slide and textbox
for slide in pptx.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for rule in formatting_rules['words']:
                    for word in rule['words']:
                        format_run(paragraph, word, rule['format'])

# Save the formatted presentation once all formatting is applied
pptx.save('06 theWhileLoop_formatted.pptx')
